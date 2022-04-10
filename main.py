import sys, os, pwd
from pptx import Presentation
from PyQt6.QtWidgets import *

version = "V 1.0 Alpha"
def get_username():
    return pwd.getpwuid(os.getuid())[0]

class filedialog(QWidget):
    def __init__(self, parent=None):
        super(filedialog, self).__init__(parent)
        layout = QVBoxLayout()
        self.btn = QPushButton("Browse")
        self.btn.clicked.connect(self.getfile)
        self.setGeometry(200, 200, 400, 200)
        self.layout2 = QHBoxLayout()
        self.deckT = QLabel("Deck:")
        self.deck = QLineEdit(self)
        self.deck.move(20, 20)
        self.deck.resize(180, 20)
        self.cardT = QLabel("Card Type (Default is Basic):")
        self.card = QLineEdit(self)
        self.card.move(20, 20)
        self.card.resize(180, 20)
        # Create a button in the window

        self.le = QLabel("")
        self.le1 = QLabel("")
        self.run = QPushButton("Run")
        self.run.setEnabled(False)
        self.run.clicked.connect(self.moveToFolder)
        layout.addWidget(self.deckT)
        layout.addWidget(self.deck)
        layout.addWidget(self.cardT)
        layout.addWidget(self.card)
        layout.addWidget(self.btn)
        layout.addWidget(self.le)
        layout.addWidget(self.le1)
        layout.addWidget(self.run)

        #layout.addWidget(self.contents)
        self.setLayout(layout)
        self.setWindowTitle("PowerPoint Anki Importer " + version)

    def getfile(self):
        uname = str(get_username())
        path = "/Users/" + uname + "/"
        fname = QFileDialog.getOpenFileName(self, 'Open file', path, "PowerPoint files (*.pptx)")
        self.le1.setText(str(fname[0]))
        self.run.setEnabled(True)
        card_type = self.card.text()
        if card_type == "":
            card_type = "Basic"
        self.le.setText(card_type + " -> " + self.deck.text())

    def show_popup(self):
        msg = QMessageBox(text="Done!", parent=self)
        msg.setIcon(QMessageBox.Icon.Information)
        msg.setStandardButtons(QMessageBox.StandardButton.Ok)
        msg.setInformativeText("The cards have been made an will be imported when you next open anki (they can also be imported manually)")
        ret = msg.exec()

    def moveToFolder(self):
        file = str(self.le1.text())
        card_type = self.card.text()
        if card_type == "":
            card_type = "Basic"
        makeCSV(file, card_type, self.deck.text())
        self.show_popup()


def makeCSV(path, c_type, deck):
    ppt=Presentation(path)
    notes = []
    uname = str(get_username())
    f_path = "/Users/" + uname + "/.anki-PPT/"
    conf = open(f_path +"conf.csv", 'a')
    conf.write(c_type + "," + deck)
    conf.write('\n')
    conf.close()
    f = open(f_path + "slides.csv", "a")
    for page, slide in enumerate(ppt.slides):
    # this is the notes that doesn't appear on the ppt slide,
    # but really the 'presenter' note.
        textNote = slide.notes_slide.notes_text_frame.text
        f.write(textNote + '\n')
        notes.append((page,textNote))
    f.close()


def main():
    app = QApplication(sys.argv)
    ex = filedialog()
    ex.show()
    sys.exit(app.exec())


if __name__ == '__main__':
    main()
